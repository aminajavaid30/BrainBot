from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from contextlib import asynccontextmanager
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.document_loaders import WebBaseLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_openai import OpenAIEmbeddings
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_openai import ChatOpenAI
from langchain_groq import ChatGroq
from langchain.chains import create_history_aware_retriever, create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_community.chat_message_histories import ChatMessageHistory
from langchain_core.chat_history import BaseChatMessageHistory
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.runnables.history import RunnableWithMessageHistory
from transformers import pipeline
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from PIL import Image
import base64
import requests
import docx2txt
import pptx
import os
import utils

## APPLICATION LIFESPAN
# Load the environment variables using FastAPI lifespan event so that they are available throughout the application
@asynccontextmanager
async def lifespan(app: FastAPI):
    # Load the environment variables
    load_dotenv()
    #os.environ['OPENAI_API_KEY'] = os.getenv("OPENAI_API_KEY")
    ## Langsmith tracking
    os.environ["LANGCHAIN_TRACING_V2"] = "true" # Enable tracing to capture all the monitoring results
    os.environ["LANGCHAIN_API_KEY"] = os.getenv("LANGCHAIN_API_KEY")
    ## load the Groq API key
    os.environ['GROQ_API_KEY'] = os.getenv("GROQ_API_KEY")
    global image_to_text
    image_to_text = pipeline("image-to-text", model="Salesforce/blip-image-captioning-large")
    yield
    # Delete all the temporary images
    utils.unlink_images("/images")

## FASTAPI APP
# Initialize the FastAPI app
app = FastAPI(lifespan=lifespan)

## PYDANTIC MODELS
# Define an APIKey Pydantic model for the request body
class APIKey(BaseModel):
    api_key: str

# Define a FileInfo Pydantic model for the request body
class FileInfo(BaseModel):
    file_path: str
    file_type: str

# Define an Image Pydantic model for the request body
class Image(BaseModel):
    image_path: str

# Define a Website Pydantic model for the request body
class Website(BaseModel):
    website_link: str

# Define a Question Pydantic model for the request body
class Question(BaseModel):
    question: str
    resource: str

## FUNCTIONS
# Function to combine all documents
def format_docs(docs):
    return "\n\n".join(doc.page_content for doc in docs)

# Function to encode the image
def encode_image(image_path):
  with open(image_path, "rb") as image_file:
    return base64.b64encode(image_file.read()).decode('utf-8')

## FASTAPI ENDPOINTS
## GET - /
@app.get("/")
async def welcome():
    return "Welcome to Brainbot!"

## POST - /set_api_key
@app.post("/set_api_key")
async def set_api_key(api_key: APIKey):
    os.environ["OPENAI_API_KEY"] = api_key.api_key
    return "API key set successfully!"

## POST - /load_file
# Load the file, split it into document chunks, and upload the document embeddings into a vectorstore   
@app.post("/load_file/{llm}")
async def load_file(llm: str, file_info: FileInfo):
    file_path = file_info.file_path
    file_type = file_info.file_type
    
    # Read the file and split it into document chunks
    try:
        # Initialize the text splitter
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

        # Check the file type and load each file according to its type
        if file_type == "application/pdf":
            # Read pdf file
            loader = PyPDFLoader(file_path)
            docs = loader.load()
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            # Read docx file
            text = docx2txt.process(file_path)
            docs = text_splitter.create_documents([text])
        elif file_type == "text/plain":
            # Read txt file
            with open(file_path, 'r') as file:
                text = file.read()
                docs = text_splitter.create_documents([text])
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            # Read pptx file
            presentation = pptx.Presentation(file_path)
            # Initialize an empty list to store slide texts
            slide_texts = []

            # Iterate through slides and extract text
            for slide in presentation.slides:
                # Initialize an empty string to store text for each slide
                slide_text = ""
            
                # Iterate through shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"  # Add shape text to slide text
                        # Append slide text to the list
                        slide_texts.append(slide_text.strip())

            docs = text_splitter.create_documents(slide_texts)
        elif file_type == "text/html":
            # Read html file
            with open(file_path, 'r') as file:
                soup = BeautifulSoup(file, 'html.parser')
                text = soup.get_text()
                docs = text_splitter.create_documents([text])

        # Delete the temporary file
        os.unlink(file_path)

        # Split the document into chunks
        documents = text_splitter.split_documents(docs)

        if llm == "GPT-4":
            embeddings = OpenAIEmbeddings()
        elif llm == "GROQ":
            embeddings = HuggingFaceEmbeddings()
        
        # Save document embeddings into the FAISS vectorstore
        global file_vectorstore
        file_vectorstore = FAISS.from_documents(documents, embeddings)
    except Exception as e:
        # Handle errors
        raise HTTPException(status_code=500, detail=str(e.with_traceback))
    return "File uploaded successfully!"

## POST - /image
# Interpret the image using the LLM - OpenAI Vision
@app.post("/image/{llm}")
async def interpret_image(llm: str, image: Image):
    try:
        # Get the base64 string
        base64_image = encode_image(image.image_path)
        
        if llm == "GPT-4":
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {os.environ['OPENAI_API_KEY']}"
            }

            payload = {
                "model": "gpt-4-turbo",
                "messages": [
                    {
                    "role": "user",
                    "content": [
                        {
                        "type": "text",
                        "text": "What's in this image?"
                        },
                        {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{base64_image}"
                        }
                        }
                    ]
                    }
                ],
                "max_tokens": 300
            }

            response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
            response = response.json()
            # Extract description about the image
            description = response["choices"][0]["message"]["content"]
        elif llm == "GROQ":
            # Use image-to-text model from Hugging Face
            response = image_to_text(image.image_path)
            # Extract description about the image
            description = response[0]["generated_text"]
            chat = ChatGroq(temperature=0, groq_api_key=os.environ["GROQ_API_KEY"], model_name="Llama3-8b-8192")
            system = "You are an assistant to understand and interpret images."
            human = "{text}"
            prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])

            chain = prompt | chat
            text = f"Explain the following image description in a small paragraph. {description}"
            response = chain.invoke({"text": text})
            description = str.capitalize(description) + ". " + response.content
    except Exception as e:
        # Handle errors
        raise HTTPException(status_code=500, detail=str(e))

    return description

## POST - load_link
# Load the website content through scraping, split it into document chunks, and upload the document
# embeddings into a vectorstore
@app.post("/load_link/{llm}")
async def website_info(llm: str, link: Website):
    try:
        # load, chunk, and index the content of the html page
        loader = WebBaseLoader(web_paths=(link.website_link,),)

        global web_documents
        web_documents = loader.load()

        # split the document into chunks
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        documents = text_splitter.split_documents(web_documents)

        if llm == "GPT-4":
            embeddings = OpenAIEmbeddings()
        elif llm == "GROQ":
            embeddings = HuggingFaceEmbeddings()

        # Save document embeddings into the FAISS vectorstore
        global website_vectorstore
        website_vectorstore = FAISS.from_documents(documents, embeddings)
    except Exception as e:
        # Handle errors
        raise HTTPException(status_code=500, detail=str(e))

    return "Website loaded successfully!"

## POST - /answer_with_chat_history
# Retrieve the answer to the question using LLM and the RAG chain maintaining the chat history
@app.post("/answer_with_chat_history/{llm}")
async def get_answer_with_chat_history(llm: str, question: Question):
    user_question = question.question
    resource = question.resource
    selected_llm = llm

    try:
        # Initialize the LLM
        if selected_llm == "GPT-4":
            llm = ChatOpenAI(model="gpt-4-turbo", temperature=0)
        elif selected_llm == "GROQ":
            llm = ChatGroq(groq_api_key=os.environ["GROQ_API_KEY"], model_name="Llama3-8b-8192")

        # extract relevant context from the document using the retriever with similarity search
        if resource == "file":
            retriever = file_vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 5})
        elif resource == "web":
            retriever = website_vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 5})
        
        ### Contextualize question ###
        contextualize_q_system_prompt = """Given a chat history and the latest user question \
        which might reference context in the chat history, formulate a standalone question \
        which can be understood without the chat history. Do NOT answer the question, \
        just reformulate it if needed and otherwise return it as is."""
        contextualize_q_prompt = ChatPromptTemplate.from_messages(
            [
                ("system", contextualize_q_system_prompt),
                MessagesPlaceholder("chat_history"),
                ("human", "{input}"),
            ]
        )
        history_aware_retriever = create_history_aware_retriever(
            llm, retriever, contextualize_q_prompt
        )

        ### Answer question ###
        qa_system_prompt = """You are an assistant for question-answering tasks. \
        Use the following pieces of retrieved context to answer the question. \
        If you don't know the answer, just say that you don't know. \
        Use three sentences maximum and keep the answer concise.\
        {context}"""
        qa_prompt = ChatPromptTemplate.from_messages(
            [
                ("system", qa_system_prompt),
                MessagesPlaceholder("chat_history"),
                ("human", "{input}"),
            ]
        )
        
        question_answer_chain = create_stuff_documents_chain(llm, qa_prompt)

        rag_chain = create_retrieval_chain(history_aware_retriever, question_answer_chain)
        
        ### Statefully manage chat history ###
        store = {}
        def get_session_history(session_id: str) -> BaseChatMessageHistory:
            if session_id not in store:
                store[session_id] = ChatMessageHistory()
            return store[session_id]
        
        conversational_rag_chain = RunnableWithMessageHistory(
            rag_chain,
            get_session_history,
            input_messages_key="input",
            history_messages_key="chat_history",
            output_messages_key="answer",
        )
        
        response = conversational_rag_chain.invoke(
            {"input": user_question},
            config={
                "configurable": {"session_id": "abc123"}
            },  # constructs a key "abc123" in `store`.
        )["answer"]
    except Exception as e:
        # Handle errors
        raise HTTPException(status_code=500, detail=str(e))
    
    return response
